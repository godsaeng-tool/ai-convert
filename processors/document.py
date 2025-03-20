import os
import fitz  # PyMuPDF
import pptx
from config import logger

def extract_text_from_pdf(file_path):
    """PDF 파일에서 텍스트 추출 - 구조 유지 방식으로 개선"""
    try:
        text_content = ""
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # 구조화된 텍스트 추출
            text_dict = page.get_text("dict")
            page_text = ""
            
            for block in text_dict.get("blocks", []):
                if block["type"] == 0:  # 텍스트 블록만 처리
                    for line in block.get("lines", []):
                        line_text = ""
                        for span in line.get("spans", []):
                            line_text += span["text"] + " "
                        page_text += line_text.strip() + "\n"
                    page_text += "\n"  # 블록 사이 줄바꿈 추가
            
            # 페이지 구분자 추가
            text_content += f"--- 페이지 {page_num+1} ---\n{page_text}\n\n"
        
        doc.close()
        logger.info(f"PDF 텍스트 추출 완료: {file_path}")
        return text_content
    except Exception as e:
        logger.error(f"PDF 텍스트 추출 실패: {str(e)}")
        return ""

def extract_text_from_ppt(file_path):
    """PPT 파일에서 텍스트 추출 - 구조 정보 유지"""
    try:
        presentation = pptx.Presentation(file_path)
        text_content = ""
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = f"--- 슬라이드 {slide_num+1} ---\n"
            
            # 제목 추출
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                slide_text += f"# {title}\n\n"
            
            # 내용 추출 (텍스트 상자, 표 등)
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 제목이 아닌 텍스트 추가
                    if shape != slide.shapes.title:
                        slide_content.append(shape.text.strip())
            
            # 슬라이드 내용 추가
            if slide_content:
                slide_text += "\n".join(slide_content) + "\n\n"
            
            text_content += slide_text
        
        logger.info(f"PPT 텍스트 추출 완료: {file_path}")
        return text_content
    except Exception as e:
        logger.error(f"PPT 텍스트 추출 실패: {str(e)}")
        return ""

def process_document(file_path, task_id, result_dir, data_dir):
    """문서 파일(PDF/PPT) 처리 및 텍스트 추출"""
    try:
        # 확장자 확인
        file_ext = file_path.split('.')[-1].lower()
        
        # 파일 타입에 따라 적절한 처리 함수 호출
        if file_ext == 'pdf':
            extracted_text = extract_text_from_pdf(file_path)
        elif file_ext in ['ppt', 'pptx']:
            extracted_text = extract_text_from_ppt(file_path)
        else:
            raise ValueError(f"지원되지 않는 문서 형식: {file_ext}")
        
        # 추출된 텍스트 저장
        text_path = os.path.join(data_dir, f"{task_id}.txt")
        with open(text_path, 'w', encoding='utf-8') as f:
            f.write(extracted_text)
        
        return {
            "success": True,
            "message": "문서 텍스트 추출 완료",
            "text_path": text_path,
            "text_content": extracted_text
        }
    
    except Exception as e:
        logger.error(f"문서 처리 실패: {str(e)}")
        return {
            "success": False,
            "message": f"문서 처리 실패: {str(e)}"
        }